import win32com.client
from PIL import ImageGrab
import os
import sys
import shutil
import datetime
import openpyxl
import pptx

rutaFotos = ''

def diccionario( book ):
    sheet = book['Encontrar Parámetros PPT']
    
    fila = 3
    columna = 2
    
    i = 0
    tamano = 0
    
    while sheet.cell( row=(i+fila), column=(columna) ).value != None:
        i += 1
    
    tamano = i
    
    dic = {}
    fots = {}
    
    for i in range( tamano ):
        dic[sheet.cell( row=(i+fila), column=(columna) ).value] = sheet.cell( row=(i+fila), column=(columna+2) ).value
        fots[sheet.cell( row=(i+fila), column=(columna) ).value] = sheet.cell( row=(i+fila), column=(columna+1) ).value
        
    return [dic, fots]
    
def guardaGraficosComoImagenes(inputExcelFilePath, hoja):

    global rutaFotos

    rutaFotos = inputExcelFilePath[:inputExcelFilePath.rindex('/')]+ '/Fotos/'

    if os.path.exists( rutaFotos ):
        outputPNGImagePath = rutaFotos
    else:
        os.makedirs( rutaFotos )
        outputPNGImagePath = rutaFotos


    # Get directory name
    mydir = win32com.__gen_path__

    if os.path.exists(mydir):
        shutil.rmtree(mydir)

    o = win32com.client.Dispatch("Excel.Application")

    wb = o.Workbooks.Open(inputExcelFilePath)


    sheet = o.Sheets(hoja)
    for n, shape in enumerate(sheet.ChartObjects()):
        

        shape.Copy()
        image = ImageGrab.grabclipboard()

        image.save(str(outputPNGImagePath+shape.Name+'.png'), 'png')
        pass
    pass

    wb.Close(True)
    o.Quit()

def ubicarPalabraP( rutaPPT, rutaExcel, me  ):
    global rutaFotos

    presentacion = pptx.Presentation(rutaPPT )
    book = openpyxl.load_workbook( filename=rutaExcel, data_only=True )    
    
    params = diccionario( book )
    
    dic = params[0]
    fots = params[1]
    
    
    guardaGraficosComoImagenes(rutaExcel, 'Encontrar Parámetros PPT')      

    for slide in presentacion.slides: 
        for shape in slide.shapes: 
            if not shape.has_text_frame: 
                continue 
            for paragraph in shape.text_frame.paragraphs: 
                for run in paragraph.runs: 
                    for val in dic:
                        if val in run.text:
                            
                            if fots[val] == "Foto":
                        
                                run.text = run.text.replace( val, "" )
                                slide.shapes.add_picture( rutaFotos + dic[val] + '.png', pptx.util.Inches(0.5), pptx.util.Inches(1.75) )
                            
                            if fots[val] != "Foto":
                                
                                if isinstance(dic[val], datetime.date):
                                    run.text = run.text.replace( val, str(dic[val].strftime('%d/%m/%y')) )
                                else:
                                    run.text = run.text.replace( val, str(dic[val]) )
                    
            
    presentacion.save( rutaPPT[:len(rutaPPT)-5] + ' ' +list(dic.values())[0] + '.pptx')
    me['text'] = 'Hecho'
    me['state'] = 'disabled'
    me['command'] = None
    