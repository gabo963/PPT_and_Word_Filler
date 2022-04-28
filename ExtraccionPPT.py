import openpyxl
import re
import pptx


def sacarVals( presentacion ):

    texto = []
    
    for slide in presentacion.slides: 
        for shape in slide.shapes: 
            if not shape.has_text_frame: 
                continue 
            for paragraph in shape.text_frame.paragraphs: 
                for run in paragraph.runs: 

                    encontrados = re.findall('\[[^\]]*\]', run.text)
    
                    for encontrado in encontrados:
                            texto.append( encontrado )

    texto = list(dict.fromkeys(texto))
    
    
    orden = []
    
    for linea in texto:
        orden.append(int(str(linea)[1:-1]))
    
    orden.sort() 
    
    texto = []
    
    for linea in orden:
        texto.append( "[" + str(linea) + "]" )
         
    return texto

def agregar_a_excelP( rutaPPT, rutaExcel, me):
    
    presentacion = pptx.Presentation(rutaPPT)
    book = openpyxl.load_workbook( rutaExcel )
    
    lista = sacarVals( presentacion )
    
    if not 'Encontrar Parámetros PPT' in book.sheetnames:
        book.create_sheet('Encontrar Parámetros PPT')
        
    sheet = book['Encontrar Parámetros PPT']
    
    fila = 3
    columna = 2
    
    i = 0
    aborrar = 0
    
    while sheet.cell( row=(i+fila), column=(columna) ).value != None:
        i += 1
    
    aborrar = i
    
    for i in range( aborrar ):
        sheet.cell( row=(i+fila), column=(columna) ).value = ""
        sheet.cell( row=(i+fila), column=(columna+1) ).value = ""
        sheet.cell( row=(i+fila), column=(columna+2) ).value = ""
        
        
    fila = 3
    columna = 2
    
    sheet.cell( row=(2), column=(2) ).value = 'Parámetro:'
    sheet.cell( row=(2), column=(3) ).value = 'Descripción:'
    sheet.cell( row=(2), column=(4) ).value = 'Valor:'
    
    for i in range( len(lista) ):
        sheet.cell( row=(i+fila), column=(columna) ).value = lista[i]
    
    book.save(rutaExcel)
    
    me['text'] = 'Hecho'
    me['state'] = 'disabled'
    me['command'] = None
    